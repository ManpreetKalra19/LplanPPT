import streamlit as st
from openai import OpenAI
import google.generativeai as genai
import os
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
import re

# Specifically for OpenAI client to work on Streamlit Cloud
import openai
openai.api_key = None  # Will be set later with user input

st.set_page_config(page_title="AI Lesson Plan Generator", page_icon="📚")

st.title("AI Lesson Plan Generator")
st.markdown("Generate custom lesson plans using AI models")

# Initialize session state
if 'lesson_plan' not in st.session_state:
    st.session_state.lesson_plan = ""
if 'create_ppt' not in st.session_state:
    st.session_state.create_ppt = False
if 'ppt_generated' not in st.session_state:
    st.session_state.ppt_generated = False
if 'ppt_path' not in st.session_state:
    st.session_state.ppt_path = None

# Create tabs for different sections
tab1, tab2 = st.tabs(["Generate Plan", "About"])

with tab1:
    # API Selection
    ai_choice = st.radio(
        "Choose which AI model to use:",
        ["OpenAI (GPT-4)", "Google Gemini"]
    )

    # API Key input with password masking
    if ai_choice == "OpenAI (GPT-4)":
        api_key = st.text_input("Enter your OpenAI API Key:", type="password")
        model_choice = st.selectbox(
            "Select OpenAI model:",
            ["gpt-4", "gpt-3.5-turbo"]
        )
    else:  # Gemini
        api_key = st.text_input("Enter your Google Gemini API Key:", type="password")
        # Using the correct model name for Gemini API
        model_choice = "gemini-1.5-pro"  # Fixed model name

    # Board selection
    board = st.radio(
        "Choose the educational board:",
        ["CBSE (Central Board of Secondary Education)", "IB (International Baccalaureate)"]
    )
    
    # Format board name
    board_name = "CBSE" if "CBSE" in board else "IB"

    # Create columns for input fields
    col1, col2 = st.columns(2)
    
    with col1:
        grade = st.text_input("Enter the Grade Level:")
        subject = st.text_input("Enter the Subject:")
    
    with col2:
        concept = st.text_input("Enter the Concept:")

    def generate_lesson_plan_openai(grade, subject, concept, board, api_key, model):
        """Generates a lesson plan using OpenAI's models with the v1.0+ API."""
        try:
            # Set API key for both approaches
            openai.api_key = api_key
            
            # Use the simpler approach that works on Streamlit Cloud
            prompt = f"""Create a detailed lesson plan for {subject} on {concept} for Grade {grade} following the {board} curriculum.
            
            Structure your response with clear headings for each section:
            1. Learning Objectives
            2. Key Concepts
            3. Required Materials
            4. Introduction Activity
            5. Main Teaching Points
            6. Student Activities
            7. Assessment Methods
            8. Homework/Extension Activities
            
            Use proper Markdown formatting with headings (##) for each section.
            """
            
            messages = [
                {"role": "system", "content": f"You are an expert educator specializing in the {board} curriculum."},
                {"role": "user", "content": prompt}
            ]
            
            response = openai.chat.completions.create(
                model=model,
                messages=messages
            )
            
            return response.choices[0].message.content
        
        except Exception as e:
            error_message = str(e)
            if "insufficient_quota" in error_message:
                return "ERROR: OpenAI API quota exceeded. Please check your account balance or try using the Gemini option instead."
            else:
                return f"OpenAI Error: {error_message}"

    def generate_lesson_plan_gemini(grade, subject, concept, board, api_key):
        """Generates a lesson plan using Google's Gemini model."""
        try:
            # Configure the API
            genai.configure(api_key=api_key)
            
            # Create a model instance
            model = genai.GenerativeModel("gemini-1.5-pro")  # Using the correct model identifier
            
            prompt = f"""
            You are an expert educator specializing in the {board} curriculum.
            Create a detailed lesson plan for {subject} on {concept} for Grade {grade}.
            
            Structure your response with clear headings for each section:
            1. Learning Objectives
            2. Key Concepts
            3. Required Materials
            4. Introduction Activity
            5. Main Teaching Points
            6. Student Activities
            7. Assessment Methods
            8. Homework/Extension Activities
            
            Use proper Markdown formatting with headings (##) for each section.
            """
            
            response = model.generate_content(prompt)
            return response.text
        
        except Exception as e:
            return f"Gemini Error: {str(e)}"

    def create_ppt_from_lesson_plan(lesson_plan, title, grade, subject, board):
        """Creates a PowerPoint presentation from the lesson plan."""
        # Create a new presentation
        prs = Presentation()
        
        # Set up slide layouts
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]
        
        # Add title slide
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = f"{subject}: {title}"
        subtitle_shape.text = f"Grade {grade} - {board} Curriculum"
        
        # Parse the lesson plan into sections
        # Look for markdown headers (## Section Name)
        sections = re.split(r'(?m)^##\s+(.*?)$', lesson_plan)
        
        # The first element is either empty or content before the first heading
        if sections[0].strip():
            # Add an overview slide if there's content before the first heading
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = "Overview"
            slide.placeholders[1].text = sections[0].strip()
        
        # Process each section (heading and content pairs)
        for i in range(1, len(sections), 2):
            if i+1 < len(sections):
                heading = sections[i].strip()
                content = sections[i+1].strip()
                
                # Add a slide for this section
                slide = prs.slides.add_slide(content_slide_layout)
                slide.shapes.title.text = heading
                
                # Format the content - remove extra whitespace and bullet points
                content = re.sub(r'\n\s*\n', '\n', content)  # Remove empty lines
                
                # Add content to slide
                text_frame = slide.placeholders[1].text_frame
                text_frame.text = ""
                
                # Process content line by line
                lines = content.split('\n')
                paragraph = text_frame.paragraphs[0]
                for line in lines:
                    # Skip empty lines
                    if not line.strip():
                        continue
                    
                    # Check if this is a bullet point
                    if line.strip().startswith('* ') or line.strip().startswith('- '):
                        if paragraph.text:  # If we already have text, add a new paragraph
                            paragraph = text_frame.add_paragraph()
                        paragraph.text = line.strip()[2:]  # Remove the bullet point marker
                        paragraph.level = 1  # Set to bullet point
                    else:
                        # Regular text
                        if paragraph.text:  # If we already have text, add a new paragraph
                            paragraph = text_frame.add_paragraph()
                        paragraph.text = line
        
        # Add final slide
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "Thank You"
        slide.placeholders[1].text = "Any questions?"
        
        # Save to a temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_file.name)
        return temp_file.name

    # Generate button
    if st.button("Generate Lesson Plan"):
        if not grade or not subject or not concept:
            st.error("Please fill in all fields (Grade, Subject, and Concept).")
        elif not api_key:
            st.error("Please enter your API key.")
        else:
            with st.spinner("Generating lesson plan..."):
                if ai_choice == "OpenAI (GPT-4)":
                    st.session_state.lesson_plan = generate_lesson_plan_openai(grade, subject, concept, board_name, api_key, model_choice)
                else:
                    st.session_state.lesson_plan = generate_lesson_plan_gemini(grade, subject, concept, board_name, api_key)

    # Display the generated lesson plan
    if st.session_state.lesson_plan:
        st.subheader("Generated Lesson Plan")
        st.markdown(st.session_state.lesson_plan)
        
        # Download button for the lesson plan
        filename = f"{board_name}_{subject}_{concept}_Grade{grade}.txt"
        st.download_button(
            label="Download Lesson Plan",
            data=st.session_state.lesson_plan,
            file_name=filename,
            mime="text/plain"
        )
        
        # Option to create PowerPoint
        if not st.session_state.create_ppt and not st.session_state.ppt_generated:
            if st.button("Create PowerPoint Presentation"):
                st.session_state.create_ppt = True
        
        # Create PowerPoint if requested
        if st.session_state.create_ppt and not st.session_state.ppt_generated:
            with st.spinner("Creating PowerPoint presentation..."):
                try:
                    ppt_path = create_ppt_from_lesson_plan(
                        st.session_state.lesson_plan,
                        concept,
                        grade,
                        subject,
                        board_name
                    )
                    st.session_state.ppt_path = ppt_path
                    st.session_state.ppt_generated = True
                    st.session_state.create_ppt = False
                except Exception as e:
                    st.error(f"Error creating PowerPoint: {str(e)}")
                    st.session_state.create_ppt = False
        
        # Download PowerPoint if generated
        if st.session_state.ppt_generated and st.session_state.ppt_path:
            with open(st.session_state.ppt_path, "rb") as ppt_file:
                ppt_bytes = ppt_file.read()
                
            st.download_button(
                label="Download PowerPoint Presentation",
                data=ppt_bytes,
                file_name=f"{board_name}_{subject}_{concept}_Grade{grade}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
            # Show success message
            st.success("PowerPoint presentation created successfully!")

with tab2:
    st.subheader("About this app")
    st.markdown("""
    This app uses AI models to generate customized lesson plans for educators:
    
    * **OpenAI GPT-4/GPT-3.5**: Advanced AI models for detailed, coherent lesson plans
    * **Google Gemini 1.5 Pro**: Google's powerful generative AI for structured lesson plans
    
    The app supports both CBSE and IB curriculums and provides output that includes learning objectives, key concepts, activities, and assessments.
    
    ### Features:
    
    * Generate detailed lesson plans with structured sections
    * Download lesson plans as text files
    * Create PowerPoint presentations from your lesson plans
    * Support for multiple educational boards
    
    **Note**: You will need your own API keys to use this application.
    
    ### Setup Instructions:
    
    #### For OpenAI:
    1. Create an account on [OpenAI Platform](https://platform.openai.com/)
    2. Generate an API key under the API Keys section
    3. Add billing information to use the API
    
    #### For Google Gemini:
    1. Visit [Google AI Studio](https://makersuite.google.com/app/apikey) or [Google AI for Developers](https://ai.google.dev/)
    2. Create or sign in to your Google account
    3. Generate an API key for Gemini API
    """)